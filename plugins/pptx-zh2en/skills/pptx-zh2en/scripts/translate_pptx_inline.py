#!/usr/bin/env python3
"""
PPTX 中英双向翻译脚本 v2.1

用法：
    extract:  python translate_pptx_inline.py --mode extract -i zh.pptx -t trans.json [--direction zh2en|en2zh]
    apply:    python translate_pptx_inline.py --mode apply   -i zh.pptx -t trans.json -o en.pptx [--direction zh2en|en2zh]

v2.1 新增：
    --direction en2zh  支持英文→中文翻译（字体/标点/质检方向反转）

v2.0 核心改进：
    1. 段落级提取+翻译（解决 run 碎片化导致翻译断裂）
    2. 中文标点自动清洗
    3. 翻译后自动质检（残留中文/标点/重复词）
    4. 双语重复检测（中英同框时提示用已有英文）
"""

import argparse
import json
import os
import re
import sys
import io
import zipfile
from collections import Counter

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx",
                           "--break-system-packages", "-q"])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE


# ═════════════════════════════════════════════════════════════════════════
#  字体映射（三层策略，v1.3 保留）
# ═════════════════════════════════════════════════════════════════════════

KNOWN_CN_FONTS = {
    "Microsoft YaHei": "Calibri", "Microsoft YaHei UI": "Calibri",
    "Microsoft YaHei Light": "Calibri Light",
    "微软雅黑": "Calibri", "微软雅黑 Light": "Calibri Light",
    "SimHei": "Arial", "黑体": "Arial", "STHeiti": "Arial", "华文黑体": "Arial",
    "Source Han Sans SC": "Arial", "Source Han Sans CN": "Arial",
    "Source Han Serif SC": "Times New Roman", "Source Han Serif CN": "Times New Roman",
    "Noto Sans CJK SC": "Arial", "Noto Serif CJK SC": "Times New Roman",
    "WenQuanYi Micro Hei": "Arial", "文泉驿微米黑": "Arial", "WenQuanYi Zen Hei": "Arial",
    "SimSun": "Times New Roman", "宋体": "Times New Roman",
    "NSimSun": "Times New Roman", "新宋体": "Times New Roman",
    "STSong": "Times New Roman", "华文宋体": "Times New Roman",
    "STZhongsong": "Times New Roman", "华文中宋": "Times New Roman",
    "FangSong": "Times New Roman", "仿宋": "Times New Roman",
    "STFangsong": "Times New Roman", "华文仿宋": "Times New Roman",
    "KaiTi": "Times New Roman", "楷体": "Times New Roman",
    "STKaiti": "Times New Roman", "华文楷体": "Times New Roman",
    "SimSun-ExtB": "Consolas", "FangSong_GB2312": "Courier New",
    "方正书宋": "Times New Roman", "方正黑体": "Arial",
    "方正仿宋": "Times New Roman", "方正楷体": "Times New Roman",
    "汉仪旗黑": "Arial", "造字工房明黑": "Arial",
}

CN_FONT_KEYWORDS = [
    "YaHei", "SimSun", "SimHei", "FangSong", "KaiTi", "STSong", "STHeiti",
    "STKaiti", "STFangsong", "NSimSun", "STZhongsong",
    "Source Han", "Noto Sans CJK", "Noto Serif CJK", "WenQuanYi",
    "MingLiU", "PMingLiU", "DFKai", "Microsoft JhengHei",
]


def is_cjk_font(font_name):
    if not font_name:
        return False
    if re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', font_name):
        return True
    fn_lower = font_name.lower()
    return any(kw.lower() in fn_lower for kw in CN_FONT_KEYWORDS)


def get_en_font_for(cn_font):
    if cn_font in KNOWN_CN_FONTS:
        return KNOWN_CN_FONTS[cn_font]
    if re.search(r'宋|Song|Ming|Serif|细明|楷|Kai', cn_font, re.IGNORECASE):
        return "Times New Roman"
    if re.search(r'黑|Hei|Sans|Gothic|粗|Bold', cn_font, re.IGNORECASE):
        return "Arial"
    if re.search(r'仿|Fang', cn_font, re.IGNORECASE):
        return "Times New Roman"
    return "Calibri"


def contains_chinese(text):
    return bool(re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', text))


def contains_english(text):
    """判断文本中是否含有英文字母"""
    return bool(re.search(r'[A-Za-z]', text))


# ─── 英文字体 → 中文字体映射（en2zh 方向使用）─────────────────────────
EN2ZH_FONT_MAP = {
    "Calibri": "微软雅黑", "Calibri Light": "微软雅黑 Light",
    "Arial": "黑体", "Arial Black": "黑体", "Arial Narrow": "黑体",
    "Times New Roman": "宋体", "Georgia": "宋体",
    "Cambria": "宋体", "Garamond": "宋体",
    "Helvetica": "黑体", "Helvetica Neue": "黑体",
    "Tahoma": "微软雅黑", "Trebuchet MS": "微软雅黑",
    "Verdana": "微软雅黑", "Segoe UI": "微软雅黑",
    "Courier New": "仿宋", "Consolas": "仿宋",
}
# en2zh 方向：中文标点替换表（英文→中文）
EN2ZH_PUNCT_MAP = {
    ', ': '，',  '. ': '。',  '; ': '；',  ': ': '：',
    '?': '？',  '!': '！',
    '(': '（',  ')': '）',
}


def clean_english_punctuation(text):
    """en2zh 方向：将英文标点替换为中文标点（仅在有中文的段落中）"""
    if not contains_chinese(text):
        return text
    # 只替换中文文本中的英文标点
    for en, zh in EN2ZH_PUNCT_MAP.items():
        text = text.replace(en, zh)
    return text


# ═════════════════════════════════════════════════════════════════════════
#  中文标点 → 英文标点
# ═════════════════════════════════════════════════════════════════════════

PUNCT_MAP = {
    '\u3001': ', ',    # 、顿号
    '\uff0c': ', ',    # ，全角逗号
    '\u3002': '. ',    # 。句号
    '\uff1b': '; ',    # ；
    '\uff1a': ': ',    # ：
    '\uff1f': '? ',    # ？
    '\uff01': '! ',    # ！
    '\uff08': '(',     # （
    '\uff09': ')',     # ）
    '\u3010': '[',     # 【
    '\u3011': ']',     # 】
    '\u201c': '"',     # "
    '\u201d': '"',     # "
    '\u2018': "'",     # '
    '\u2019': "'",     # '
    '\u300a': '<',     # 《
    '\u300b': '>',     # 》
}


def clean_chinese_punctuation(text):
    """将残留的中文标点替换为英文标点"""
    for cn, en in PUNCT_MAP.items():
        text = text.replace(cn, en)
    # 修正连续标点：", ," → ",", "//" → "/"
    text = re.sub(r',\s*,', ',', text)
    text = re.sub(r'/\s*/', '/', text)
    # 修正 "? ?" → "?"
    text = re.sub(r'\?\s*\?', '?', text)
    # 修正句末标点前的多余空格
    text = re.sub(r'\s+([.,;:!?])', r'\1', text)
    return text


# ═════════════════════════════════════════════════════════════════════════
#  字体替换（三层）
# ═════════════════════════════════════════════════════════════════════════

def replace_run_font(run, was_translated=False):
    try:
        if not run.font:
            return False
        fn = run.font.name
        if fn:
            if fn in KNOWN_CN_FONTS:
                run.font.name = KNOWN_CN_FONTS[fn]
                return True
            if is_cjk_font(fn):
                run.font.name = get_en_font_for(fn)
                return True
            return False
        if fn is None and was_translated:
            run.font.name = "Calibri"
            return True
    except Exception:
        pass
    return False


def replace_run_font_en2zh(run):
    """en2zh 方向：将英文字体替换为中文字体"""
    try:
        if not run.font:
            return False
        fn = run.font.name
        if fn and fn in EN2ZH_FONT_MAP:
            run.font.name = EN2ZH_FONT_MAP[fn]
            return True
        # 如果 run.font.name 为 None，设为微软雅黑
        if fn is None:
            run.font.name = "微软雅黑"
            return True
    except Exception:
        pass
    return False


def fix_theme_fonts(pptx_path):
    buf = io.BytesIO()
    count = 0
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if 'theme' in item and item.endswith('.xml'):
                    try:
                        xml = data.decode('utf-8')
                        new_xml, cnt = _fix_theme_xml(xml)
                        if cnt > 0:
                            count += cnt
                            data = new_xml.encode('utf-8')
                    except UnicodeDecodeError:
                        pass
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())
    return count


def _fix_theme_xml(xml):
    count = 0
    def _replace_ea(m):
        nonlocal count
        old_tf = m.group(1)
        if old_tf and is_cjk_font(old_tf):
            count += 1
            return '<a:ea typeface=""/>'
        return m.group(0)
    xml = re.sub(r'<a:ea\s+typeface="([^"]*)"\s*/>', _replace_ea, xml)
    for script in ["Hans", "Hant", "Jpan", "Hang"]:
        def _replace_script(m, s=script):
            nonlocal count
            old_tf = m.group(1)
            if old_tf and is_cjk_font(old_tf):
                count += 1
                return f'<a:font script="{s}" typeface="Calibri"/>'
            return m.group(0)
        xml = re.sub(
            rf'<a:font\s+script="{script}"\s+typeface="([^"]*)"\s*/>',
            _replace_script, xml)
    return xml, count


def fix_theme_fonts_en2zh(pptx_path):
    """en2zh 方向：将主题字体中的英文字体替换为中文字体"""
    buf = io.BytesIO()
    count = 0
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if 'theme' in item and item.endswith('.xml'):
                    try:
                        xml = data.decode('utf-8')
                        # 替换 <a:ea typeface=""/> 为 <a:ea typeface="微软雅黑"/>
                        def _rep_ea(m):
                            nonlocal count
                            count += 1
                            return '<a:ea typeface="微软雅黑"/>'
                        xml = re.sub(r'<a:ea\s+typeface=""\s*/>', _rep_ea, xml)
                        # 替换 Hans/Hant 字体
                        for script in ["Hans", "Hant"]:
                            def _rep_sc(m, s=script):
                                nonlocal count
                                count += 1
                                return f'<a:font script="{s}" typeface="微软雅黑"/>'
                            xml = re.sub(
                                rf'<a:font\s+script="{script}"\s+typeface="([^"]*)"\s*/>',
                                _rep_sc, xml)
                        data = xml.encode('utf-8')
                    except UnicodeDecodeError:
                        pass
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())
    return count


# ═════════════════════════════════════════════════════════════════════════
#  递归遍历形状
# ═════════════════════════════════════════════════════════════════════════

def iter_shapes_recursive(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from iter_shapes_recursive(child)
        return
    if shape.has_table:
        for r, row in enumerate(shape.table.rows):
            for c, cell in enumerate(row.cells):
                yield shape, cell.text_frame, f"table[{r}][{c}]"
        return
    if shape.has_text_frame:
        yield shape, shape.text_frame, "body"


def iter_slide_text_frames(slide):
    for shape in slide.shapes:
        yield from iter_shapes_recursive(shape)
    if slide.has_notes_slide:
        for shape in slide.notes_slide.shapes:
            if shape.has_text_frame:
                yield shape, shape.text_frame, "notes"


# ═════════════════════════════════════════════════════════════════════════
#  双语重复检测：提取段落中已有的英文部分
# ═════════════════════════════════════════════════════════════════════════

def extract_existing_english(para_text):
    """
    检测段落中是否已有英文翻译。
    很多中文 PPT 在中文后附带了英文，如 "语义记忆 Semantic Memory"
    返回提取到的英文部分，如果没有则返回空字符串。
    """
    # 匹配：中文字符 + 空格 + 英文单词（至少2个连续英文词）
    m = re.search(r'[\u4e00-\u9fff\u3400-\u4dbf][\s]*([A-Za-z][A-Za-z\s/\-&()]+[A-Za-z)])\s*$', para_text)
    if m:
        en = m.group(1).strip()
        # 至少要有2个词才认为是翻译
        if len(en.split()) >= 2:
            return en
    # 也匹配：英文在前，中文在后
    m = re.search(r'^([A-Za-z][A-Za-z\s/\-&()]+[A-Za-z)])\s*[\u4e00-\u9fff]', para_text)
    if m:
        en = m.group(1).strip()
        if len(en.split()) >= 2:
            return en
    return ""


# ═════════════════════════════════════════════════════════════════════════
#  extract 模式（段落级）
# ═════════════════════════════════════════════════════════════════════════

def extract_texts(pptx_path, output_json, direction="zh2en"):
    prs = Presentation(pptx_path)
    items = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        for shape, tf, sub_path in iter_slide_text_frames(slide):
            location = "notes" if sub_path == "notes" else "slide"
            for para_idx, para in enumerate(tf.paragraphs):
                # 段落级提取：拼接所有 run 的文本
                runs = para.runs
                if not runs:
                    continue
                full_text = "".join(r.text for r in runs)
                if not full_text.strip():
                    continue

                # 根据方向决定提取哪些段落
                if direction == "zh2en":
                    if not contains_chinese(full_text):
                        continue
                    existing_en = extract_existing_english(full_text)
                    existing_translation = existing_en
                    has_existing = bool(existing_en)
                else:  # en2zh
                    if not contains_english(full_text):
                        continue
                    # 跳过纯数字/符号
                    if not re.search(r'[A-Za-z]{2,}', full_text):
                        continue
                    existing_translation = ""
                    has_existing = False

                # 获取首个 run 的字体名
                first_font = ""
                if runs[0].font and runs[0].font.name:
                    first_font = runs[0].font.name

                items.append({
                    "id": len(items),
                    "slide": slide_num,
                    "shape_name": shape.name,
                    "shape_id": shape.shape_id,
                    "sub_path": sub_path,
                    "para_idx": para_idx,
                    "run_count": len(runs),
                    "location": location,
                    "original": full_text,
                    "translated": "",
                    "font_name": first_font,
                    "has_existing_english": has_existing,
                    "existing_english": existing_translation,
                })

    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(items, f, ensure_ascii=False, indent=2)

    print(f"✅ 提取完成：共 {len(items)} 个段落（段落级提取）")
    print(f"   已保存至: {output_json}")
    if items:
        locs = Counter(it["sub_path"] for it in items)
        detail = ", ".join(f"{k}={v}" for k, v in locs.most_common())
        print(f"   来源分布: {detail}")
        en_count = sum(1 for it in items if it["has_existing_english"])
        if en_count:
            print(f"   ⚡ {en_count} 个段落已含英文翻译（JSON 中 existing_english 字段）")
    print(f"\n请翻译 JSON 中所有 'translated' 字段，然后运行 --mode apply")
    print(f"💡 提示：has_existing_english=true 的条目，可直接将 existing_english 填入 translated")


# ═════════════════════════════════════════════════════════════════════════
#  文本框自动适配（防止翻译后溢出）
# ═════════════════════════════════════════════════════════════════════════

# 最小字号（pt），不允许缩到比这更小
MIN_FONT_SIZE = 8


def auto_fit_text_frame(shape, tf, direction="zh2en"):
    """
    对翻译后的文本框做自动适配，防止文字溢出。
    策略：
      - 如果原文已有 spAutoFit（形状自适应）→ 不动，保持原样
      - 如果原文已有 normAutofit → 不动，PowerPoint 会自动缩放
      - 如果原文是 noAutofit 或无设置 → 开启 word_wrap + 设置 normAutofit
      - 不主动缩小字号（避免字号过小）
      - 不修改原有的 word_wrap 设置（False → True 会改变排版）
    """
    from pptx.oxml.ns import qn

    try:
        body_pr = tf._txBody.find(qn('a:bodyPr'))
        if body_pr is None:
            return

        # 检查现有 autofit 类型
        has_sp_autofit = body_pr.find(qn('a:spAutoFit')) is not None
        has_norm_autofit = body_pr.find(qn('a:normAutofit')) is not None
        has_no_autofit = body_pr.find(qn('a:noAutofit')) is not None

        # 如果已经有 spAutoFit 或 normAutofit，完全不动
        if has_sp_autofit or has_norm_autofit:
            return

        # 只有原来明确是 noAutofit 的才改为 normAutofit
        # 如果原来没有任何 autofit 设置（如 Rounded Rectangle），不动它
        if has_no_autofit:
            body_pr.remove(body_pr.find(qn('a:noAutofit')))
            norm_autofit = body_pr.makeelement(qn('a:normAutofit'), {})
            body_pr.append(norm_autofit)

    except Exception:
        pass


# ═════════════════════════════════════════════════════════════════════════
#  apply 模式（段落级写入 + 标点清洗 + 质检）
# ═════════════════════════════════════════════════════════════════════════

def apply_translations(pptx_path, translations_json, output_path, direction="zh2en"):
    with open(translations_json, "r", encoding="utf-8") as f:
        items = json.load(f)

    missing = [it for it in items if not it.get("translated", "").strip()]
    if missing:
        print(f"⚠️  警告：有 {len(missing)} 个段落未翻译（translated 为空），将保留原文")

    # 检测"翻译等于原文"（可能遗漏翻译）
    unchanged = []
    for it in items:
        tr = it.get("translated", "").strip()
        orig = it.get("original", "").strip()
        if tr and tr == orig:
            # 判断是否真的需要翻译（根据方向检测源语言字符）
            if direction == "zh2en" and contains_chinese(orig):
                unchanged.append(it)
            elif direction == "en2zh" and contains_english(orig) and re.search(r'[A-Z][a-z]{3,}', orig):
                # 排除学术引用和文件路径
                if 'et al.' not in orig and not orig.startswith('file:') and not orig.startswith('/') and '|' not in orig:
                    unchanged.append(it)
    if unchanged:
        print(f"⚠️  警告：有 {len(unchanged)} 个段落 translated == original（疑似未翻译）：")
        for it in unchanged[:10]:
            print(f"   [幻灯片{it['slide']}] {it['original'][:60]}")
        if len(unchanged) > 10:
            print(f"   ... 还有 {len(unchanged)-10} 个（省略）")

    prs = Presentation(pptx_path)
    changed = 0
    font_run_changed = 0
    warnings = []

    # 索引：(slide, shape_id, sub_path, para_idx, location) → translated
    tmap = {}
    for it in items:
        key = (it["slide"], it["shape_id"], it["sub_path"],
               it["para_idx"], it["location"])
        tmap[key] = it.get("translated", "") or it["original"]

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        for shape, tf, sub_path in iter_slide_text_frames(slide):
            location = "notes" if sub_path == "notes" else "slide"

            # 先收集该 shape 的所有段落文本，检测是否已有英文翻译
            # 例如 "工作记忆\nWorking Memory" → 只保留英文
            all_paras = list(tf.paragraphs)
            full_shape_text = "\n".join(p.text for p in all_paras)

            # 双语去重：如果 shape 内同时有中文和英文段落
            # 且英文段落是对中文段落的翻译，则只保留英文
            if direction == "zh2en" and contains_chinese(full_shape_text):
                en_paras = []
                zh_paras = []
                for p in all_paras:
                    ptext = p.text.strip()
                    if not ptext:
                        continue
                    if contains_chinese(ptext) and not re.search(r'[A-Z][a-z]{2,}', ptext):
                        zh_paras.append(p)
                    elif re.search(r'[A-Z][a-z]{2,}', ptext) and not contains_chinese(ptext):
                        en_paras.append(p)

                # 如果中文段落数 == 英文段落数 > 0，说明是中英对照
                if zh_paras and en_paras and len(zh_paras) == len(en_paras):
                    # 只翻译中文段落，英文段落保持不变
                    for p in zh_paras:
                        para_idx = all_paras.index(p)
                        key = (slide_num, shape.shape_id, sub_path,
                               para_idx, location)
                        if key in tmap:
                            translated = tmap[key]
                            if translated and translated != p.text:
                                if direction == "zh2en":
                                    translated = clean_chinese_punctuation(translated)
                                runs = p.runs
                                if runs:
                                    runs[0].text = ""
                                    for r in runs[1:]:
                                        r.text = ""
                                    changed += 1
                                    for r in runs:
                                        if replace_run_font(r, was_translated=True):
                                            font_run_changed += 1
                    auto_fit_text_frame(shape, tf, direction)
                    continue

            # 正常处理（非双语对照的情况）
            for para_idx, para in enumerate(all_paras):
                key = (slide_num, shape.shape_id, sub_path,
                       para_idx, location)
                if key not in tmap:
                    continue

                runs = para.runs
                if not runs:
                    continue

                original = "".join(r.text for r in runs)
                translated = tmap[key]
                if not translated or translated == original:
                    continue

                # 标点清洗（根据方向）
                if direction == "zh2en":
                    translated = clean_chinese_punctuation(translated)
                else:  # en2zh
                    translated = clean_english_punctuation(translated)

                # 溢出检测
                if direction == "zh2en":
                    src_w = sum(2 if '\u4e00' <= c <= '\u9fff' else 1 for c in original)
                else:
                    src_w = len(original)
                ratio = len(translated) / max(src_w, 1)
                if ratio > 1.5:
                    warnings.append(
                        f"幻灯片{slide_num} '{shape.name}': "
                        f"视觉宽度增加约 {ratio:.1f}x"
                    )

                # 写入第一个 run，清空其余 run
                runs[0].text = translated
                for r in runs[1:]:
                    r.text = ""
                changed += 1

                # 字体替换（根据方向）
                for r in runs:
                    if direction == "zh2en":
                        if replace_run_font(r, was_translated=True):
                            font_run_changed += 1
                    else:  # en2zh：英文字体 → 中文字体
                        if replace_run_font_en2zh(r):
                            font_run_changed += 1

            # 文本框级适配
            auto_fit_text_frame(shape, tf, direction)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    # 主题字体替换
    if direction == "zh2en":
        font_theme_changed = fix_theme_fonts(output_path)
    else:
        font_theme_changed = fix_theme_fonts_en2zh(output_path)

    target_lang = "中文版" if direction == "en2zh" else "英文版"
    print(f"✅ 文本写入完成：已替换 {changed} 个段落")
    if font_run_changed:
        print(f"✅ Run 级字体替换：{font_run_changed} 处")
    if font_theme_changed:
        print(f"✅ 主题字体替换：{font_theme_changed} 处")
    print(f"   {target_lang}已保存至: {output_path}")

    if warnings:
        print(f"\n⚠️  可能溢出的文本框（共 {len(warnings)} 处）：")
        for w in warnings[:10]:
            print(f"   {w}")

    # 自动质检
    qa_results = auto_qa(output_path, direction)
    if qa_results["issues"]:
        print(f"\n🔍 自动质检发现 {len(qa_results['issues'])} 个问题：")
        for issue in qa_results["issues"][:15]:
            print(f"   {issue}")
        if len(qa_results["issues"]) > 15:
            print(f"   ... 还有 {len(qa_results['issues'])-15} 个（省略）")
    else:
        if direction == "zh2en":
            print(f"\n✅ 自动质检通过：无残留中文、无标点问题、无重复词")
        else:
            print(f"\n✅ 自动质检通过：无残留英文、无标点问题")


# ═════════════════════════════════════════════════════════════════════════
#  自动质检
# ═════════════════════════════════════════════════════════════════════════

def auto_qa(pptx_path, direction="zh2en"):
    """翻译后自动质检，返回问题列表"""
    prs = Presentation(pptx_path)
    issues = []

    all_puncts = set(PUNCT_MAP.keys())

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        for shape, tf, sub_path in iter_slide_text_frames(slide):
            for para_idx, para in enumerate(tf.paragraphs):
                text = para.text
                if not text.strip():
                    continue

                if direction == "zh2en":
                    # 检测残留中文
                    if contains_chinese(text):
                        issues.append(f"幻灯片{slide_num} '{shape.name}': 残留中文 → {text[:40]}")
                    # 检测残留中文标点
                    found_puncts = [p for p in all_puncts if p in text]
                    if found_puncts:
                        issues.append(f"幻灯片{slide_num} '{shape.name}': 残留中文标点 {found_puncts} → {text[:40]}")
                    # 重复单词
                    dup_match = re.search(r'\b(\w{2,})\s+\1\b', text, re.IGNORECASE)
                    if dup_match:
                        issues.append(f"幻灯片{slide_num} '{shape.name}': 疑似重复词 '{dup_match.group(1)}' → {text[:40]}")
                    # 缺失空格
                    space_match = re.search(r'[a-z][A-Z][a-z]', text)
                    if space_match:
                        if not re.match(r'^[a-z]+[A-Z][a-z]+$', text.split()[0] if text.split() else ''):
                            issues.append(f"幻灯片{slide_num} '{shape.name}': 疑似缺失空格 '{space_match.group()}' → {text[:40]}")
                else:  # en2zh
                    # 检测残留英文单词（排除缩写、公式、引用）
                    # 跳过：单独的字母（如 C3, C4）、公式、作者引用
                    en_words = re.findall(r'\b[A-Z][a-z]{3,}\b', text)
                    # 过滤掉常见不需要翻译的英文
                    skip = {"PAR", "RGB", "LiDAR", "SOP", "RUE", "LAI", "NDVI",
                            "UAV", "C3", "C4", "API", "AI", "DNA", "RNA", "GPS"}
                    real_en = [w for w in en_words if w not in skip and not re.match(r'^[A-Z]+$', w)]
                    if real_en:
                        issues.append(f"幻灯片{slide_num} '{shape.name}': 疑似残留英文 {real_en[:3]} → {text[:40]}")

    return {"issues": issues}


# ═════════════════════════════════════════════════════════════════════════
#  主入口
# ═════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="PPTX 中英双向翻译工具 v2.1")
    parser.add_argument("--mode", choices=["extract", "apply"], required=True)
    parser.add_argument("--direction", choices=["zh2en", "en2zh"], default="zh2en",
                        help="翻译方向：zh2en（中→英，默认）| en2zh（英→中）")
    parser.add_argument("--input", "-i", required=True, help="输入 PPTX 路径")
    parser.add_argument("--output", "-o", help="输出 PPTX 路径（apply 模式）")
    parser.add_argument("--translations-file", "-t", help="翻译 JSON 路径")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"❌ 输入文件不存在: {args.input}")
        sys.exit(1)

    if args.mode == "extract":
        json_path = args.translations_file or args.input.replace(".pptx", "_translations.json")
        extract_texts(args.input, json_path, args.direction)
    elif args.mode == "apply":
        if not args.translations_file or not os.path.isfile(args.translations_file):
            print(f"❌ 翻译文件不存在: {args.translations_file}")
            sys.exit(1)
        output = args.output or args.input.replace(".pptx", "_ZH.pptx" if args.direction == "en2zh" else "_EN.pptx")
        apply_translations(args.input, args.translations_file, output, args.direction)


if __name__ == "__main__":
    main()
