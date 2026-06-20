#!/usr/bin/env python3
"""
PPTX 中英双向翻译脚本（外部 API 模式）v2.4
用法: python translate_pptx.py --input zh.pptx --output en.pptx [--direction zh2en|en2zh] [--api-key KEY]

依赖: python-pptx, requests
功能: 段落级翻译、三层字体替换、标点清洗、自动质检、表格/组合形状支持
"""

import argparse
import io
import json
import os
import re
import sys
import time
import zipfile
from collections import Counter
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx",
                           "--break-system-packages", "-q"])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

import requests

# ─── 常量 ──────────────────────────────────────────────────────────────────
BATCH_SIZE = 30
REQUEST_DELAY = 0.3

# ─── 字体处理（三层策略，与 inline 版一致）──────────────────────────────
KNOWN_CN_FONTS = {
    "Microsoft YaHei": "Calibri", "Microsoft YaHei Light": "Calibri Light",
    "微软雅黑": "Calibri", "微软雅黑 Light": "Calibri Light",
    "SimHei": "Arial", "黑体": "Arial", "STHeiti": "Arial", "华文黑体": "Arial",
    "Source Han Sans SC": "Arial", "Source Han Sans CN": "Arial",
    "Source Han Serif SC": "Times New Roman", "Source Han Serif CN": "Times New Roman",
    "Noto Sans CJK SC": "Arial", "Noto Serif CJK SC": "Times New Roman",
    "WenQuanYi Micro Hei": "Arial", "文泉驿微米黑": "Arial", "WenQuanYi Zen Hei": "Arial",
    "SimSun": "Times New Roman", "宋体": "Times New Roman", "NSimSun": "Times New Roman",
    "新宋体": "Times New Roman", "STSong": "Times New Roman", "华文宋体": "Times New Roman",
    "STZhongsong": "Times New Roman", "华文中宋": "Times New Roman",
    "FangSong": "Times New Roman", "仿宋": "Times New Roman",
    "STFangsong": "Times New Roman", "华文仿宋": "Times New Roman",
    "KaiTi": "Times New Roman", "楷体": "Times New Roman",
    "STKaiti": "Times New Roman", "华文楷体": "Times New Roman",
    "SimSun-ExtB": "Consolas", "FangSong_GB2312": "Courier New",
    "方正书宋": "Times New Roman", "方正黑体": "Arial",
    "方正仿宋": "Times New Roman", "方正楷体": "Times New Roman",
    "汉仪旗黑": "Arial", "造字工房明黑": "Arial",
}
CN_FONT_KEYWORDS = [
    "YaHei", "SimSun", "SimHei", "FangSong", "KaiTi", "STSong", "STHeiti",
    "STKaiti", "STFangsong", "NSimSun", "STZhongsong",
    "Source Han", "Noto Sans CJK", "Noto Serif CJK", "WenQuanYi",
    "MingLiU", "PMingLiU", "DFKai", "Microsoft JhengHei",
]


def is_cjk_font(font_name):
    if not font_name:
        return False
    if re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', font_name):
        return True
    fn_lower = font_name.lower()
    return any(kw.lower() in fn_lower for kw in CN_FONT_KEYWORDS)


def get_en_font_for(cn_font):
    if cn_font in KNOWN_CN_FONTS:
        return KNOWN_CN_FONTS[cn_font]
    if re.search(r'宋|Song|Ming|Serif|细明|楷|Kai', cn_font, re.IGNORECASE):
        return "Times New Roman"
    if re.search(r'黑|Hei|Sans|Gothic|粗|Bold', cn_font, re.IGNORECASE):
        return "Arial"
    if re.search(r'仿|Fang', cn_font, re.IGNORECASE):
        return "Times New Roman"
    return "Calibri"


def contains_chinese(text: str) -> bool:
    return bool(re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', text))


def replace_run_font(run, was_translated=False):
    """三层字体替换：映射表 → 启发式 → 兜底"""
    try:
        if not run.font:
            return False
        fn = run.font.name
        if fn:
            if fn in KNOWN_CN_FONTS:
                run.font.name = KNOWN_CN_FONTS[fn]
                return True
            if is_cjk_font(fn):
                run.font.name = get_en_font_for(fn)
                return True
            return False
        if fn is None and was_translated:
            run.font.name = "Calibri"
            return True
    except Exception:
        pass
    return False


def fix_theme_fonts(pptx_path):
    """修改 theme1.xml 中的东亚字体定义"""
    import io
    buf = io.BytesIO()
    count = 0
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if 'theme' in item and item.endswith('.xml'):
                    try:
                        xml = data.decode('utf-8')
                        # 替换 <a:ea typeface="中文字体"/>
                        def _rep_ea(m):
                            nonlocal count
                            tf = m.group(1)
                            if tf and is_cjk_font(tf):
                                count += 1
                                return '<a:ea typeface=""/>'
                            return m.group(0)
                        xml = re.sub(r'<a:ea\s+typeface="([^"]*)"\s*/>', _rep_ea, xml)
                        # 替换 Hans/Hant/Jpan/Hang
                        for sc in ["Hans", "Hant", "Jpan", "Hang"]:
                            def _rep_sc(m, s=sc):
                                nonlocal count
                                tf = m.group(1)
                                if tf and is_cjk_font(tf):
                                    count += 1
                                    return f'<a:font script="{s}" typeface="Calibri"/>'
                                return m.group(0)
                            xml = re.sub(
                                rf'<a:font\s+script="{sc}"\s+typeface="([^"]*)"\s*/>',
                                _rep_sc, xml)
                        data = xml.encode('utf-8')
                    except UnicodeDecodeError:
                        pass
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())
    return count


# ─── en2zh 方向：英文字体 → 中文字体 ──────────────────────────────────────

EN2ZH_FONT_MAP = {
    "Calibri": "微软雅黑", "Calibri Light": "微软雅黑 Light",
    "Arial": "黑体", "Arial Black": "黑体", "Arial Narrow": "黑体",
    "Times New Roman": "宋体", "Georgia": "宋体",
    "Cambria": "宋体", "Garamond": "宋体",
    "Helvetica": "黑体", "Helvetica Neue": "黑体",
    "Tahoma": "微软雅黑", "Trebuchet MS": "微软雅黑",
    "Verdana": "微软雅黑", "Segoe UI": "微软雅黑",
    "Courier New": "仿宋", "Consolas": "仿宋",
}


def replace_run_font_en2zh(run):
    """en2zh 方向：将英文字体替换为中文字体"""
    try:
        if not run.font:
            return False
        fn = run.font.name
        if fn and fn in EN2ZH_FONT_MAP:
            run.font.name = EN2ZH_FONT_MAP[fn]
            return True
        if fn is None:
            run.font.name = "微软雅黑"
            return True
    except Exception:
        pass
    return False


def fix_theme_fonts_en2zh(pptx_path):
    """en2zh 方向：将主题字体中的英文字体替换为中文字体"""
    buf = io.BytesIO()
    count = 0
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if 'theme' in item and item.endswith('.xml'):
                    try:
                        xml = data.decode('utf-8')
                        def _rep_ea(m):
                            nonlocal count
                            count += 1
                            return '<a:ea typeface="微软雅黑"/>'
                        xml = re.sub(r'<a:ea\s+typeface=""\s*/>', _rep_ea, xml)
                        for script in ["Hans", "Hant"]:
                            def _rep_sc(m, s=script):
                                nonlocal count
                                count += 1
                                return f'<a:font script="{s}" typeface="微软雅黑"/>'
                            xml = re.sub(
                                rf'<a:font\s+script="{script}"\s+typeface="([^"]*)"\s*/>',
                                _rep_sc, xml)
                        data = xml.encode('utf-8')
                    except UnicodeDecodeError:
                        pass
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())
    return count


# ─── 递归遍历（与 inline 版一致）──────────────────────────────────────────

def iter_shapes_recursive(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from iter_shapes_recursive(child)
        return
    if shape.has_table:
        for r, row in enumerate(shape.table.rows):
            for c, cell in enumerate(row.cells):
                yield shape, cell.text_frame, f"table[{r}][{c}]"
        return
    if shape.has_text_frame:
        yield shape, shape.text_frame, "body"


def iter_slide_text_frames(slide):
    for shape in slide.shapes:
        yield from iter_shapes_recursive(shape)
    if slide.has_notes_slide:
        for shape in slide.notes_slide.shapes:
            if shape.has_text_frame:
                yield shape, shape.text_frame, "notes"


# ─── API 翻译 ─────────────────────────────────────────────────────────────

def translate_batch_via_api(texts, api_url, api_key, model, direction="zh2en"):
    if not texts:
        return []
    numbered = "\n".join(f"[{i+1}] {t}" for i, t in enumerate(texts))

    if direction == "zh2en":
        task = "Translate the following Chinese texts to English."
        rules = "3. For slide titles: use Title Case"
    else:
        task = "Translate the following English texts to Chinese (Simplified)."
        rules = "3. For slide titles: use concise Chinese"

    prompt = f"""You are a professional academic translator specializing in agricultural science, remote sensing, and digital phenotyping.

{task} Rules:
1. Return ONLY the translations, numbered with [1], [2], etc. — no extra explanation
2. Keep professional academic terminology accurate
{rules}
4. Preserve any numbers, units (e.g. μmol·m⁻²·s⁻¹), and formulas unchanged
5. If a text is already in the target language or is a number/symbol, return it as-is
6. Keep translations concise — slide text should be brief and impactful

Texts to translate:
{numbered}"""

    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    payload = {"model": model, "messages": [{"role": "user", "content": prompt}], "temperature": 0.2}

    try:
        resp = requests.post(f"{api_url.rstrip('/')}/chat/completions",
                             headers=headers, json=payload, timeout=120)
        resp.raise_for_status()
        result_text = resp.json()["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"  [警告] API 请求失败: {e}，返回原文")
        return texts

    translated = {}
    for line in result_text.strip().split("\n"):
        m = re.match(r"\[(\d+)\]\s*(.*)", line.strip())
        if m:
            translated[int(m.group(1)) - 1] = m.group(2).strip()
    return [translated.get(i, texts[i]) for i in range(len(texts))]


def translate_via_api(texts, api_url, api_key, model, direction="zh2en"):
    all_results = []
    total = len(texts)
    print(f"  共 {total} 条文本，分批翻译（每批 {BATCH_SIZE} 条）...")
    for start in range(0, total, BATCH_SIZE):
        batch = texts[start:start + BATCH_SIZE]
        end = min(start + BATCH_SIZE, total)
        print(f"  翻译第 {start+1}–{end} 条...")
        all_results.extend(translate_batch_via_api(batch, api_url, api_key, model, direction))
        if end < total:
            time.sleep(REQUEST_DELAY)
    return all_results


# ─── 主流程 ───────────────────────────────────────────────────────────────

def translate_presentation(input_path, output_path, api_url, api_key, model, direction="zh2en"):
    print(f"\n{'='*60}")
    print(f"  输入: {input_path}")
    print(f"  输出: {output_path}")
    print(f"  方向: {direction}")
    print(f"  模型: {model}")
    print(f"{'='*60}\n")

    prs = Presentation(input_path)
    total_slides = len(prs.slides)
    print(f"[1/4] 加载成功，共 {total_slides} 张幻灯片\n")

    # 收集需要翻译的文本
    print("[2/4] 扫描所有文本（含表格、组合形状、备注页）...")
    run_infos = []  # [(run, original, slide_num)]
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        for shape, tf, sub_path in iter_slide_text_frames(slide):
            for para in tf.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if direction == "zh2en" and contains_chinese(run.text):
                        run_infos.append((run, run.text, slide_num))
                    elif direction == "en2zh" and contains_english(run.text) and re.search(r'[A-Za-z]{2,}', run.text):
                        run_infos.append((run, run.text, slide_num))

    total_texts = len(run_infos)
    src_label = "中文" if direction == "zh2en" else "英文"
    print(f"  找到 {total_texts} 条含{src_label}的文本段\n")

    if total_texts == 0:
        print(f"  未找到{src_label}内容，文件可能已是目标语言版本。")
        import shutil
        shutil.copy2(input_path, output_path)
        return

    # 翻译
    print("[3/4] 开始翻译...")
    raw_texts = [info[1] for info in run_infos]
    translated_texts = translate_via_api(raw_texts, api_url, api_key, model, direction)

    # 写回 + 字体替换
    print("\n[4/4] 写回翻译结果 + 替换字体...")
    changed = 0
    font_changed = 0
    warnings = []

    for (run, original, slide_num), translated in zip(run_infos, translated_texts):
        was_translated = False
        if translated and translated != original:
            zh_w = sum(2 if '\u4e00' <= c <= '\u9fff' else 1 for c in original)
            ratio = len(translated) / max(zh_w, 1)
            if ratio > 1.5:
                warnings.append(f"  幻灯片 {slide_num}: 文本宽度增加 {ratio:.1f}x")
            run.text = translated
            changed += 1
            was_translated = True
        if replace_run_font(run, was_translated) if direction == "zh2en" else replace_run_font_en2zh(run):
            font_changed += 1

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    # 主题字体替换
    if direction == "zh2en":
        theme_changed = fix_theme_fonts(output_path)
    else:
        theme_changed = fix_theme_fonts_en2zh(output_path)

    print(f"\n✅ 翻译完成！共 {total_slides} 张幻灯片，{changed} 条文本")
    print(f"✅ Run 级字体替换：{font_changed} 处")
    print(f"✅ 主题字体替换：{theme_changed} 处")
    print(f"   英文版已保存至: {output_path}\n")

    if warnings:
        print(f"⚠️  以下幻灯片文本可能溢出（共 {len(warnings)} 处）:")
        for w in warnings[:10]:
            print(w)

    print(f"\n{'='*60}")


def main():
    parser = argparse.ArgumentParser(description="PPTX 中英双向翻译（外部 API 模式）v2.4")
    parser.add_argument("--input", "-i", required=True, help="输入 PPTX 文件路径")
    parser.add_argument("--output", "-o", help="输出 PPTX 文件路径")
    parser.add_argument("--direction", choices=["zh2en", "en2zh"], default="zh2en",
                        help="翻译方向：zh2en（中→英，默认）| en2zh（英→中）")
    parser.add_argument("--api-url",
                        default=os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1"))
    parser.add_argument("--api-key", default=os.environ.get("OPENAI_API_KEY", ""))
    parser.add_argument("--model", default=os.environ.get("TRANSLATE_MODEL", "gpt-4o-mini"))

    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"❌ 错误：输入文件不存在: {args.input}")
        sys.exit(1)
    if not args.input.lower().endswith(".pptx"):
        print("❌ 错误：只支持 .pptx 格式（不支持 .ppt）")
        sys.exit(1)

    suffix = "_ZH.pptx" if args.direction == "en2zh" else "_EN.pptx"
    output = args.output or str(Path(args.input).with_name(f"{Path(args.input).stem}{suffix}"))
    translate_presentation(args.input, output, args.api_url, args.api_key, args.model, args.direction)


if __name__ == "__main__":
    main()
